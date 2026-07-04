"""
Generates a PowerPoint executive summary deck for the GRC / Oracle ITGC-ITAC audit:
  GRC_ITAC_Executive_Summary.pptx
Includes a risk heat map, findings visuals, SoD/access summary, monitoring snapshot
and a phased remediation timeline. For education / template use.
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ---------- Palette ----------
NAVY = RGBColor(0x1F, 0x38, 0x64)
BLUE = RGBColor(0x2E, 0x54, 0x96)
LIGHTBLUE = RGBColor(0xD9, 0xE1, 0xF2)
GREEN = RGBColor(0x2E, 0x7D, 0x32)
LGREEN = RGBColor(0xC6, 0xEF, 0xCE)
AMBER = RGBColor(0xF5, 0x9E, 0x0B)
LAMBER = RGBColor(0xFF, 0xEB, 0x9C)
RED = RGBColor(0xC0, 0x00, 0x00)
LRED = RGBColor(0xFF, 0xC7, 0xCE)
CRIT = RGBColor(0xE0, 0x00, 0x00)
GREY = RGBColor(0x59, 0x59, 0x59)
LGREY = RGBColor(0xF2, 0xF2, 0xF2)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
DARK = RGBColor(0x26, 0x26, 0x26)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]
SW, SH = prs.slide_width, prs.slide_height


def slide():
    return prs.slides.add_slide(BLANK)


def rect(s, x, y, w, h, fill, line=None, shape=MSO_SHAPE.RECTANGLE):
    sp = s.shapes.add_shape(shape, Inches(x), Inches(y), Inches(w), Inches(h))
    sp.fill.solid(); sp.fill.fore_color.rgb = fill
    if line is None:
        sp.line.fill.background()
    else:
        sp.line.color.rgb = line; sp.line.width = Pt(1)
    sp.shadow.inherit = False
    return sp


def text(s, x, y, w, h, txt, size=14, color=DARK, bold=False, align=PP_ALIGN.LEFT,
         anchor=MSO_ANCHOR.TOP, italic=False, font="Calibri"):
    tb = s.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = True; tf.vertical_anchor = anchor
    lines = txt.split("\n")
    for i, ln in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        r = p.add_run(); r.text = ln
        r.font.size = Pt(size); r.font.bold = bold; r.font.italic = italic
        r.font.color.rgb = color; r.font.name = font
    return tb


def shape_text(sp, txt, size=12, color=WHITE, bold=True, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE):
    tf = sp.text_frame; tf.word_wrap = True; tf.vertical_anchor = anchor
    for i, ln in enumerate(txt.split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        r = p.add_run(); r.text = ln
        r.font.size = Pt(size); r.font.bold = bold; r.font.color.rgb = color; r.font.name = "Calibri"


def header(s, title, kicker="GRC / ITGC-ITAC AUDIT"):
    rect(s, 0, 0, 13.333, 1.15, NAVY)
    rect(s, 0, 1.15, 13.333, 0.06, AMBER)
    text(s, 0.5, 0.12, 11, 0.35, kicker, size=11, color=LIGHTBLUE, bold=True)
    text(s, 0.5, 0.42, 12.3, 0.7, title, size=24, color=WHITE, bold=True)


def footer(s, page):
    text(s, 0.5, 7.05, 8, 0.35, "ABC Manufacturing Ltd.  |  FY 2025-26  |  Confidential - Template",
         size=8, color=GREY)
    text(s, 12.2, 7.05, 0.8, 0.35, str(page), size=9, color=GREY, align=PP_ALIGN.RIGHT)


def chip(s, x, y, w, h, label, value, fill, valcolor=WHITE):
    rect(s, x, y, w, h, fill)
    text(s, x, y + 0.12, w, 0.7, value, size=30, color=valcolor, bold=True, align=PP_ALIGN.CENTER)
    text(s, x, y + h - 0.55, w, 0.5, label, size=11, color=valcolor, align=PP_ALIGN.CENTER)


# =========================================================
# SLIDE 1: TITLE
# =========================================================
s = slide()
rect(s, 0, 0, 13.333, 7.5, NAVY)
rect(s, 0, 4.55, 13.333, 0.07, AMBER)
text(s, 0.9, 2.3, 11.5, 1.2, "ITGC & ITAC Audit", size=46, color=WHITE, bold=True)
text(s, 0.9, 3.35, 11.5, 0.8, "Executive Summary - Oracle Risk Management Cloud", size=22, color=LIGHTBLUE)
text(s, 0.9, 4.8, 11.5, 0.5, "ABC Manufacturing Ltd.   |   FY 2025-26", size=16, color=WHITE, bold=True)
text(s, 0.9, 5.25, 11.5, 0.5, "Oracle Fusion Cloud ERP (Release 24C)", size=13, color=LIGHTBLUE)
text(s, 0.9, 6.5, 11.5, 0.5, "Prepared for Audit Committee  -  Confidential  -  Template for educational use",
     size=10, color=GREY, italic=True)

# =========================================================
# SLIDE 2: EXECUTIVE SUMMARY
# =========================================================
s = slide()
header(s, "Executive Summary")
# KPI chips
chip(s, 0.5, 1.5, 2.9, 1.5, "Applications in scope", "6", BLUE)
chip(s, 3.6, 1.5, 2.9, 1.5, "Controls tested", "48", BLUE)
chip(s, 6.7, 1.5, 2.9, 1.5, "SoD conflicts (users)", "64", AMBER)
chip(s, 9.8, 1.5, 3.0, 1.5, "Material weakness", "1", RED)
# Narrative
text(s, 0.5, 3.35, 12.3, 0.4, "Overall Conclusion", size=16, color=NAVY, bold=True)
text(s, 0.5, 3.8, 12.3, 2.4,
     "The Oracle Fusion control environment is broadly functioning, with strong automated application "
     "controls (e.g., 3-way match). However, access governance requires urgent attention:\n"
     "-  1 Material Weakness: three finance users hold User Administration together with transaction posting.\n"
     "-  1 Significant Deficiency: developer access to migrate changes to production is not segregated.\n"
     "-  64 of 342 users (19%) carry at least one Segregation of Duties conflict.\n"
     "-  Automated controls tested effective; ITGC change management gap undermines reliance.\n\n"
     "A phased remediation roadmap is proposed, targeting critical access removal within 30 days.",
     size=13, color=DARK)
footer(s, 2)

# =========================================================
# SLIDE 3: SCOPE & APPROACH
# =========================================================
s = slide()
header(s, "Scope & Approach")
text(s, 0.5, 1.4, 6, 0.4, "In-Scope Applications", size=15, color=NAVY, bold=True)
apps = ["Oracle General Ledger", "Oracle Payables", "Oracle Receivables",
        "Oracle Fixed Assets", "Oracle Inventory", "Oracle HCM Payroll"]
y = 1.9
for a in apps:
    rect(s, 0.5, y, 0.16, 0.16, AMBER)
    text(s, 0.8, y - 0.08, 5.5, 0.35, a, size=12, color=DARK)
    y += 0.42
text(s, 7.0, 1.4, 6, 0.4, "Domains Tested", size=15, color=NAVY, bold=True)
doms = [("ITGC", "Access, Change, Operations, Backup"),
        ("ITAC", "Matching, calculations, validations, workflows"),
        ("SoD / Access", "87-rule ruleset via Oracle AAC (path-based)"),
        ("Transaction Monitoring", "15 ATC models (duplicate pay, ghost vendor)"),
        ("Substantive", "Reperformance / recalculation (limited)")]
y = 1.9
for k, v in doms:
    rect(s, 7.0, y, 2.0, 0.42, LIGHTBLUE)
    shape_text_box = text(s, 7.1, y + 0.02, 1.9, 0.38, k, size=11, color=NAVY, bold=True, anchor=MSO_ANCHOR.MIDDLE)
    text(s, 9.15, y + 0.02, 4.0, 0.42, v, size=11, color=DARK, anchor=MSO_ANCHOR.MIDDLE)
    y += 0.55
text(s, 0.5, 6.2, 12.3, 0.6,
     "Approach: controls-reliance strategy with limited substantive testing. Automated controls tested via "
     "test-of-one + ITGC reliance; manual/IT-dependent via attribute sampling. Overall materiality INR 12.5 Cr.",
     size=11, color=GREY, italic=True)
footer(s, 3)

# =========================================================
# SLIDE 4: RATING SUMMARY
# =========================================================
s = slide()
header(s, "Findings - Rating Summary")
data = [("Material Weakness", "1", RED, "Reasonable possibility of material misstatement not prevented/detected"),
        ("Significant Deficiency", "1", AMBER, "Merits governance attention; not material"),
        ("Control Deficiency", "2", RGBColor(0xBF, 0x8F, 0x00), "Control missing or not operating as designed"),
        ("Observations", "3", BLUE, "Process improvement opportunities")]
x = 0.5
for label, val, col, desc in data:
    rect(s, x, 1.5, 3.0, 1.4, col)
    text(s, x, 1.62, 3.0, 0.8, val, size=40, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    text(s, x, 2.5, 3.0, 0.35, label, size=12, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    text(s, x, 3.05, 3.0, 1.3, desc, size=10.5, color=DARK, align=PP_ALIGN.CENTER)
    x += 3.16
text(s, 0.5, 4.7, 12.3, 0.4, "Severity Assessment Logic", size=14, color=NAVY, bold=True)
text(s, 0.5, 5.15, 12.3, 1.4,
     "Each deficiency is evaluated on: (1) Could it cause a misstatement?  (2) Magnitude?  "
     "(3) Likelihood?  (4) Are there effective compensating controls?  ->  classified as Control "
     "Deficiency, Significant Deficiency, or Material Weakness.",
     size=12.5, color=DARK)
footer(s, 4)

# =========================================================
# SLIDE 5: RISK HEAT MAP
# =========================================================
s = slide()
header(s, "Risk Heat Map - Likelihood vs Impact")
# grid origin
ox, oy = 2.3, 1.7
cell = 1.35
rows = 3  # likelihood: High(top), Med, Low(bottom)
cols = 3  # impact: Low, Med, High
# color matrix by (row_from_top, col): risk increases toward top-right
colormatrix = [
    [LAMBER, AMBER, LRED],   # High likelihood
    [LGREEN, LAMBER, AMBER],  # Med
    [LGREEN, LGREEN, LAMBER], # Low
]
for r in range(rows):
    for c in range(cols):
        rect(s, ox + c * cell, oy + r * cell, cell, cell, colormatrix[r][c], line=WHITE)
# axis labels
text(s, ox - 2.0, oy, 1.9, cell, "High", size=12, color=DARK, bold=True, align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.MIDDLE)
text(s, ox - 2.0, oy + cell, 1.9, cell, "Medium", size=12, color=DARK, bold=True, align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.MIDDLE)
text(s, ox - 2.0, oy + 2 * cell, 1.9, cell, "Low", size=12, color=DARK, bold=True, align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.MIDDLE)
text(s, ox - 2.3, oy - 0.5, 2.0, 0.4, "LIKELIHOOD", size=11, color=NAVY, bold=True)
text(s, ox, oy + 3 * cell + 0.05, cell, 0.35, "Low", size=12, color=DARK, bold=True, align=PP_ALIGN.CENTER)
text(s, ox + cell, oy + 3 * cell + 0.05, cell, 0.35, "Medium", size=12, color=DARK, bold=True, align=PP_ALIGN.CENTER)
text(s, ox + 2 * cell, oy + 3 * cell + 0.05, cell, 0.35, "High", size=12, color=DARK, bold=True, align=PP_ALIGN.CENTER)
text(s, ox + cell, oy + 3 * cell + 0.45, cell, 0.35, "IMPACT", size=11, color=NAVY, bold=True, align=PP_ALIGN.CENTER)

# plot findings as dots
def dot(s, cx, cy, label, col):
    d = 0.5
    sp = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(cx - d/2), Inches(cy - d/2), Inches(d), Inches(d))
    sp.fill.solid(); sp.fill.fore_color.rgb = col; sp.line.color.rgb = WHITE; sp.line.width = Pt(1.5)
    sp.shadow.inherit = False
    shape_text(sp, label, size=11, color=WHITE, bold=True)

# D-08: High impact, High likelihood -> top-right cell center
dot(s, ox + 2 * cell + cell/2, oy + cell/2, "D-08", CRIT)
# D-05: Med impact, Med likelihood -> center cell
dot(s, ox + cell + cell/2 + 0.35, oy + cell + cell/2, "D-05", RED)
# D-03: Low impact, High likelihood -> top-left
dot(s, ox + cell/2, oy + cell/2 + 0.35, "D-03", AMBER)
# D-11: Low impact, Low likelihood -> bottom-left
dot(s, ox + cell/2, oy + 2 * cell + cell/2, "D-11", RGBColor(0xBF, 0x8F, 0x00))

# legend
lx = ox + 3 * cell + 0.9
text(s, lx, 1.6, 4.2, 0.4, "Findings Plotted", size=14, color=NAVY, bold=True)
legend = [("D-08", "Privileged access / SoD (Material Weakness)", CRIT),
          ("D-05", "Change segregation (Significant Deficiency)", RED),
          ("D-03", "Provisioning approval (Control Deficiency)", AMBER),
          ("D-11", "Depreciation rounding (Control Deficiency)", RGBColor(0xBF, 0x8F, 0x00))]
y = 2.15
for code, desc, col in legend:
    sp = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(lx), Inches(y), Inches(0.3), Inches(0.3))
    sp.fill.solid(); sp.fill.fore_color.rgb = col; sp.line.fill.background(); sp.shadow.inherit = False
    text(s, lx + 0.42, y - 0.06, 4.0, 0.5, code + " - " + desc, size=11, color=DARK)
    y += 0.62
footer(s, 5)

# =========================================================
# SLIDE 6: KEY FINDINGS DETAIL
# =========================================================
s = slide()
header(s, "Key Findings - Detail")
cards = [
    ("D-08  |  MATERIAL WEAKNESS", RED,
     "Excessive privileged access - SoD violation",
     "3 finance users hold User Administration + transaction posting.",
     "Risk of self-granted access & fraudulent entries; potential material misstatement.",
     "Remove User Admin duty; quarterly access certification.  Target: 15-Jul-2026"),
    ("D-05  |  SIGNIFICANT DEFICIENCY", AMBER,
     "Inadequate change management segregation",
     "Developer had production migration access.",
     "Unauthorized / untested changes could reach production.",
     "Segregate dev & migration; enforce approval gate.  Target: 31-Aug-2026"),
]
y = 1.45
for tag, col, title, cond, effect, rec in cards:
    rect(s, 0.5, y, 12.3, 2.55, LGREY)
    rect(s, 0.5, y, 0.16, 2.55, col)
    text(s, 0.8, y + 0.12, 11.8, 0.35, tag, size=12, color=col, bold=True)
    text(s, 0.8, y + 0.5, 11.8, 0.4, title, size=16, color=NAVY, bold=True)
    text(s, 0.8, y + 1.0, 11.8, 0.4, "Condition: " + cond, size=12, color=DARK)
    text(s, 0.8, y + 1.42, 11.8, 0.5, "Effect: " + effect, size=12, color=DARK)
    text(s, 0.8, y + 1.98, 11.8, 0.5, "Action: " + rec, size=12, color=GREEN, bold=True)
    y += 2.75
footer(s, 6)

# =========================================================
# SLIDE 7: SoD / ACCESS SUMMARY
# =========================================================
s = slide()
header(s, "Segregation of Duties - Access Analysis")
chip(s, 0.5, 1.5, 3.0, 1.4, "Users scanned", "342", BLUE)
chip(s, 3.66, 1.5, 3.0, 1.4, "Users w/ conflict", "64", AMBER)
chip(s, 6.82, 1.5, 3.0, 1.4, "Critical conflicts", "3", RED)
chip(s, 9.98, 1.5, 2.85, 1.4, "SoD rules", "87", NAVY)
text(s, 0.5, 3.2, 12.3, 0.4, "Conflicts by Process (bar)", size=14, color=NAVY, bold=True)
bars = [("P2P", 23, RED), ("R2R", 14, AMBER), ("O2C", 10, AMBER),
        ("Payroll", 8, BLUE), ("Security", 5, RED), ("Others", 4, GREY)]
maxv = 23
bx, by, bw = 0.5, 3.7, 11.5
for i, (lbl, val, col) in enumerate(bars):
    y = by + i * 0.5
    text(s, bx, y, 1.3, 0.4, lbl, size=12, color=DARK, bold=True, anchor=MSO_ANCHOR.MIDDLE)
    w = (val / maxv) * (bw - 2.0)
    rect(s, bx + 1.4, y + 0.05, w, 0.32, col)
    text(s, bx + 1.5 + w, y, 1.5, 0.4, str(val), size=12, color=DARK, bold=True, anchor=MSO_ANCHOR.MIDDLE)
footer(s, 7)

# =========================================================
# SLIDE 8: MONITORING SNAPSHOT
# =========================================================
s = slide()
header(s, "Continuous Monitoring - Dashboard Snapshot")
kpis = [("Open incidents", "64", "v from 78", AMBER),
        ("Mean time to remediate", "28d", "target < 21d", AMBER),
        ("False-positive rate", "11%", "v from 24%", GREEN),
        ("Control coverage", "83%", "target 100%", AMBER),
        ("Overdue remediations", "4", "v from 9", AMBER),
        ("Critical open", "3", "target 0", RED)]
x, y = 0.5, 1.5
for i, (lbl, val, sub, col) in enumerate(kpis):
    cx = 0.5 + (i % 3) * 4.15
    cy = 1.5 + (i // 3) * 1.6
    rect(s, cx, cy, 3.95, 1.4, WHITE, line=col)
    rect(s, cx, cy, 0.14, 1.4, col)
    text(s, cx + 0.3, cy + 0.12, 3.5, 0.35, lbl, size=12, color=GREY, bold=True)
    text(s, cx + 0.3, cy + 0.45, 3.5, 0.6, val, size=30, color=NAVY, bold=True)
    text(s, cx + 0.3, cy + 1.02, 3.5, 0.3, sub, size=10, color=col, bold=True)
text(s, 0.5, 4.9, 12.3, 0.4, "ATC Transaction Monitoring - Value at Risk", size=14, color=NAVY, bold=True)
text(s, 0.5, 5.35, 12.3, 1.4,
     "-  Duplicate payments: 12 flagged, 2 confirmed (INR 8.4 L) - under recovery\n"
     "-  Ghost vendor (bank match): 3 flagged, 1 confirmed (INR 12.0 L) - investigating\n"
     "-  Post-termination payments: 5 flagged, 3 confirmed (INR 2.15 L) - recovered",
     size=13, color=DARK)
footer(s, 8)

# =========================================================
# SLIDE 9: REMEDIATION ROADMAP (TIMELINE)
# =========================================================
s = slide()
header(s, "Remediation Roadmap - Phased Plan")
# timeline base
ty = 2.4
rect(s, 0.8, ty + 0.7, 11.7, 0.08, NAVY)
phases = [
    ("Phase 1", "0-30 days", "Critical access removal, disable super-users, obvious SoD fixes", RED),
    ("Phase 2", "1-3 months", "Custom role redesign, AAC model deployment, workflow fixes", AMBER),
    ("Phase 3", "3-6 months", "Continuous monitoring, access certification, training", GREEN),
]
x = 0.8
w = 3.9
for i, (ph, when, desc, col) in enumerate(phases):
    cx = 0.8 + i * 3.95
    # node
    sp = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(cx + 1.7), Inches(ty + 0.5), Inches(0.5), Inches(0.5))
    sp.fill.solid(); sp.fill.fore_color.rgb = col; sp.line.color.rgb = WHITE; sp.line.width = Pt(2); sp.shadow.inherit = False
    shape_text(sp, str(i + 1), size=14, color=WHITE)
    # card above
    rect(s, cx, ty - 1.4, 3.7, 1.7, LGREY)
    rect(s, cx, ty - 1.4, 3.7, 0.45, col)
    text(s, cx, ty - 1.38, 3.7, 0.45, ph + "  (" + when + ")", size=13, color=WHITE, bold=True, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    text(s, cx + 0.15, ty - 0.85, 3.4, 1.1, desc, size=11.5, color=DARK)
text(s, 0.8, ty + 1.5, 12, 0.4, "Owners assigned per action; each remediation is re-tested and evidence re-performed before closure.",
     size=12, color=GREY, italic=True)
# quick-win callout
rect(s, 0.8, 5.6, 11.7, 1.0, LRED)
text(s, 1.0, 5.72, 11.3, 0.8,
     "Immediate priority (Phase 1): Remove User Administration duty from 3 finance users and validate via "
     "AAC re-scan (0 conflicts) by 15-Jul-2026.", size=13, color=RED, bold=True, anchor=MSO_ANCHOR.MIDDLE)
footer(s, 9)

# =========================================================
# SLIDE 10: NEXT STEPS
# =========================================================
s = slide()
header(s, "Recommendations & Next Steps")
recs = [
    ("1", "Remediate the Material Weakness immediately", "Remove privileged access from finance users; validate via AAC."),
    ("2", "Segregate change management duties", "Split developer and production-migration access; enforce approval gate."),
    ("3", "Deploy continuous SoD monitoring", "Operationalise the 87-rule AAC ruleset with scheduled synchronisation."),
    ("4", "Institute periodic access certification", "Quarterly manager review-and-approve of team access."),
    ("5", "Strengthen automated controls reliance", "Lock configuration change management to sustain ITAC reliance."),
    ("6", "Track remediation to closure", "Owner + target date per finding; re-test and evidence before closing."),
]
y = 1.55
for num, title, desc in recs:
    sp = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.6), Inches(y), Inches(0.55), Inches(0.55))
    sp.fill.solid(); sp.fill.fore_color.rgb = NAVY; sp.line.fill.background(); sp.shadow.inherit = False
    shape_text(sp, num, size=16, color=WHITE)
    text(s, 1.4, y - 0.05, 11.4, 0.4, title, size=14, color=NAVY, bold=True)
    text(s, 1.4, y + 0.32, 11.4, 0.4, desc, size=11.5, color=DARK)
    y += 0.9
footer(s, 10)

# =========================================================
# SLIDE 11: CLOSING
# =========================================================
s = slide()
rect(s, 0, 0, 13.333, 7.5, NAVY)
rect(s, 0, 3.7, 13.333, 0.07, AMBER)
text(s, 0.9, 2.7, 11.5, 0.9, "Questions & Discussion", size=40, color=WHITE, bold=True)
text(s, 0.9, 3.9, 11.5, 0.6, "Full workpaper package: 21-tab Excel workbook + supporting Word deliverables",
     size=15, color=LIGHTBLUE)
text(s, 0.9, 4.4, 11.5, 0.6, "Refer to the grc-workpapers folder for detailed testing evidence and the remediation tracker.",
     size=13, color=LIGHTBLUE)

out = "/projects/sandbox/Vivek-Kriplani/grc-workpapers/GRC_ITAC_Executive_Summary.pptx"
prs.save(out)
print("Saved:", out)
print("Slides:", len(prs.slides._sldIdLst))
